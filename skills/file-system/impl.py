import os
import json
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader
from core.env_utils import ensure_package_installed
from core.interaction import ask_user

# Lazy import helpers
def get_openpyxl():
    ensure_package_installed("openpyxl")
    import openpyxl
    return openpyxl

def _is_god_mode(context):
    if context and 'config_manager' in context:
        return context['config_manager'].get_god_mode()
    return False

def _validate_path(workspace_dir, path, context=None, must_exist=False):
    """
    Validate path security and existence.
    Returns absolute path if valid, raises exception otherwise.
    """
    if not workspace_dir:
        raise ValueError("Workspace not selected.")
    
    abs_path = os.path.abspath(os.path.join(workspace_dir, path))
    abs_workspace = os.path.abspath(workspace_dir)
    
    god_mode = _is_god_mode(context)
    
    if not god_mode and not abs_path.startswith(abs_workspace):
         raise PermissionError("Access denied (Path Traversal).")
    
    if must_exist and not os.path.exists(abs_path):
        raise FileNotFoundError(f"Path '{path}' does not exist.")
        
    return abs_path

def list_files(workspace_dir, path=".", _context=None):
    """
    List files in the current workspace directory.
    
    Args:
        workspace_dir (str): The root workspace directory (injected by system).
        path (str): Relative path to list, default is '.'.
    """
    try:
        abs_path = _validate_path(workspace_dir, path, _context, must_exist=True)
        
        items = os.listdir(abs_path)
        # Filter hidden files
        items = [i for i in items if not i.startswith('.')]
        return json.dumps(items)
    except Exception as e:
        return f"Error: {str(e)}"

def rename_file(workspace_dir, old_path, new_path, _context=None):
    """
    Rename a file or directory.
    
    Args:
        workspace_dir (str): The root workspace directory (injected by system).
        old_path (str): The current relative path of the file/directory.
        new_path (str): The new relative path of the file/directory.
    """
    try:
        # Validate source
        abs_old_path = _validate_path(workspace_dir, old_path, _context, must_exist=True)
        # Validate dest (don't require existence, but check security)
        abs_new_path = _validate_path(workspace_dir, new_path, _context, must_exist=False)
        
        if os.path.exists(abs_new_path):
            return f"Error: Destination '{new_path}' already exists."

        os.rename(abs_old_path, abs_new_path)
        return f"Success: Renamed '{old_path}' to '{new_path}'."
            
    except Exception as e:
        return f"Error: {str(e)}"

def read_file(workspace_dir, path, _context=None):
    """
    Read the content of a file. Automatically detects file type based on extension.
    
    Args:
        workspace_dir (str): The root workspace directory (injected by system).
        path (str): Relative path to the file.
    """
    try:
        # Check extension first to dispatch to specific readers
        _, ext = os.path.splitext(path)
        ext = ext.lower()
        
        if ext == '.docx':
            return read_docx(workspace_dir, path, _context)
        elif ext == '.pptx':
            return read_pptx(workspace_dir, path, _context)
        elif ext == '.xlsx':
            return read_excel(workspace_dir, path, None, _context)
        elif ext == '.pdf':
            return read_pdf(workspace_dir, path, _context)

        abs_path = _validate_path(workspace_dir, path, _context, must_exist=True)
        
        if not os.path.isfile(abs_path):
            return f"Error: '{path}' is not a file."
            
        with open(abs_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read(100 * 1024)
            
    except Exception as e:
        return f"Error: {str(e)}"

def delete_file(workspace_dir, path, _context=None):
    """
    Delete a file or empty directory.
    
    Args:
        workspace_dir (str): The root workspace directory (injected by system).
        path (str): Relative path to the file.
    """
    try:
        abs_path = _validate_path(workspace_dir, path, _context, must_exist=True)
            
        # Ask for confirmation
        # Strict check for True (Yes button). Any text response or False counts as cancellation for safety.
        if ask_user(f"⚠️ DANGER: Are you sure you want to delete '{path}'?") is not True:
            return "Error: Deletion cancelled by user."

        if os.path.isfile(abs_path):
            os.remove(abs_path)
        elif os.path.isdir(abs_path):
            os.rmdir(abs_path) # Only empty directories
        else:
            return f"Error: Unknown file type for '{path}'."
            
        return f"Success: Deleted '{path}'."
            
    except Exception as e:
        return f"Error: {str(e)}"

# --- Office Suite Functions ---

def read_docx(workspace_dir, path, _context=None):
    """
    Read text content from a DOCX file.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the DOCX file.
    """
    try:
        abs_path = _validate_path(workspace_dir, path, _context, must_exist=True)
            
        doc = Document(abs_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        return f"Error reading DOCX: {str(e)}"

def write_docx(workspace_dir, path, content, mode='w', _context=None):
    """
    Write content to a DOCX file.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the DOCX file.
        content (str): Text content to write.
        mode (str): 'w' to overwrite/create, 'a' to append.
    """
    try:
        abs_path = _validate_path(workspace_dir, path, _context, must_exist=False)
        
        if mode == 'a' and os.path.exists(abs_path):
            doc = Document(abs_path)
        else:
            doc = Document()
            
        # Split by newlines and add as paragraphs
        for line in content.split('\n'):
            doc.add_paragraph(line)
            
        doc.save(abs_path)
        return f"Success: Written to '{path}'."
    except Exception as e:
        return f"Error writing DOCX: {str(e)}"

def read_pptx(workspace_dir, path, _context=None):
    """
    Read text content from a PPTX file.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the PPTX file.
    """
    try:
        abs_path = _validate_path(workspace_dir, path, _context, must_exist=True)
            
        prs = Presentation(abs_path)
        text_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            text_content.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
            
        return "\n\n".join(text_content)
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"

def create_pptx(workspace_dir, path, slides_data, _context=None):
    """
    Create a PPTX file with given slides.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the output PPTX file.
        slides_data (list): List of dicts, each with 'title' and 'content'.
                            Example: [{"title": "Slide 1", "content": "Hello World"}]
    """
    try:
        abs_path = _validate_path(workspace_dir, path, _context, must_exist=False)
        prs = Presentation()
        
        # Ensure slides_data is a list
        if isinstance(slides_data, str):
            try:
                slides_data = json.loads(slides_data)
            except:
                return "Error: slides_data must be a JSON list or valid list object."

        for slide_info in slides_data:
            title_text = slide_info.get('title', '')
            content_text = slide_info.get('content', '')
            
            # Use a standard layout (Title and Content)
            slide_layout = prs.slide_layouts[1] 
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = title_text
            content.text = content_text
            
        prs.save(abs_path)
        return f"Success: Created presentation at '{path}'."
    except Exception as e:
        return f"Error creating PPTX: {str(e)}"

def read_excel(workspace_dir, path, sheet_name=None, _context=None):
    """
    Read data from an Excel file.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the XLSX file.
        sheet_name (str): Optional sheet name to read.
    """
    try:
        abs_path = _validate_path(workspace_dir, path, _context, must_exist=True)
            
        # Use openpyxl for lightweight reading
        openpyxl = get_openpyxl()
        wb = openpyxl.load_workbook(abs_path, data_only=True)
        
        if sheet_name:
            if sheet_name not in wb.sheetnames:
                 return f"Error: Sheet '{sheet_name}' not found. Available: {wb.sheetnames}"
            sheet = wb[sheet_name]
        else:
            sheet = wb.active
            
        rows = []
        for row in sheet.iter_rows(values_only=True):
            # Convert None to empty string for better display
            cleaned_row = [str(cell) if cell is not None else "" for cell in row]
            rows.append("\t".join(cleaned_row))
            
        return "\n".join(rows)
    except Exception as e:
        return f"Error reading Excel: {str(e)}"

def write_excel(workspace_dir, path, data, sheet_name='Sheet1', _context=None):
    """
    Write data to an Excel file.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the XLSX file.
        data (list): List of lists representing rows.
        sheet_name (str): Name of the sheet.
    """
    try:
        abs_path = _validate_path(workspace_dir, path, _context, must_exist=False)
        
        # Ensure data is a list
        if isinstance(data, str):
            try:
                data = json.loads(data)
            except:
                return "Error: data must be a JSON list of lists."
        
        openpyxl = get_openpyxl()
        
        # Check if file exists to append or create
        if os.path.exists(abs_path):
             wb = openpyxl.load_workbook(abs_path)
             if sheet_name in wb.sheetnames:
                 # If sheet exists, maybe we should clear it or append?
                 # For simplicity, let's create a new sheet if it exists or overwrite?
                 # Let's remove the old sheet and create new one to match 'overwrite' behavior
                 del wb[sheet_name]
             ws = wb.create_sheet(sheet_name)
        else:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = sheet_name
        
        for row in data:
            ws.append(row)
            
        wb.save(abs_path)
        return f"Success: Written to '{path}'."
    except Exception as e:
        return f"Error writing Excel: {str(e)}"

def read_pdf(workspace_dir, path, _context=None):
    """
    Read text from a PDF file.
    
    Args:
        workspace_dir (str): Root workspace directory.
        path (str): Relative path to the PDF file.
    """
    try:
        abs_path = _validate_path(workspace_dir, path, _context, must_exist=True)
            
        reader = PdfReader(abs_path)
        text_content = []
        
        for i, page in enumerate(reader.pages):
            text_content.append(f"--- Page {i+1} ---\n" + page.extract_text())
            
        return "\n".join(text_content)
    except Exception as e:
        return f"Error reading PDF: {str(e)}"
